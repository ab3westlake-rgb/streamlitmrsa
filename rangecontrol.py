
# app.py
import io
import itertools
import textwrap
from typing import List, Dict, Tuple
from datetime import datetime, date

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib
matplotlib.use("Agg")  # headless-safe backend
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from matplotlib.colors import TwoSlopeNorm
from statsmodels.stats.proportion import proportions_ztest, proportion_confint

# Optional export libraries (handled gracefully if missing)
pptx_available = True
pdf_available = True
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except Exception:
    pptx_available = False

try:
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader  # for drawing images if needed
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
    )
    from reportlab.lib.styles import getSampleStyleSheet
except Exception:
    pdf_available = False


# ------------------------ Streamlit UI ------------------------
st.set_page_config(page_title="MRSA Infection Rate Visual", layout="wide")
st.title("MRSA Infection Rate (per 1,000 BDOC) — Time-Range Groups")

st.markdown(
    """
    Upload a dataset containing **either**:
    - **Pre‑aggregated groups** with columns: **Group**, **Group Infections**, **Group BedDays**, optional **Alpha Selected**  
    - **Raw time‑series** with columns: **Date**, **Infections**, **BedDays** (use the *Custom date ranges* mode)

    Adjust α and display options in the sidebar.
    """
)

# Sidebar controls
with st.sidebar:
    st.header("Controls")
    alpha_override = st.checkbox("Override α (alpha) manually", value=False)
    alpha_manual = st.slider("α (significance level)", 0.001, 0.20, 0.05, 0.001)
    show_annotations = st.checkbox("Show numeric annotations in heatmap", value=True)
    show_ci_labels = st.checkbox("Show CI labels above bars", value=True)
    # Toggle to choose input style
    use_date_picker = st.checkbox(
        "Use date pickers (instead of typing dates)",
        value=True,
        help="Uncheck to enter dates as text (YYYY-MM-DD)"
    )

# >>> Mode selector ABOVE the upload <<<
mode = st.radio(
    "Select data mode",
    ["Pre‑aggregated groups", "Custom date ranges (up to 4 groups)"],
    horizontal=True,
)

# --- Uploader key (used to reset/clear files) ---
if "upload_key" not in st.session_state:
    st.session_state["upload_key"] = 0

# ----- Upload widget -----
uploaded = st.file_uploader(
    "Upload Excel or CSV file",
    type=["xlsx", "csv"],
    key=st.session_state["upload_key"]
)
st.session_state["uploaded_file"] = uploaded  # store reference so we can clear

# Sample data (used if no file uploaded) — pre-aggregated
def sample_aggregated():
    return pd.DataFrame({
        "Group": ["Range 1", "Range 2", "Range 3", "Range 4"],
        "Group Infections": [3, 12, 7, 0],
        "Group BedDays": [2500, 4800, 3100, 2100],
        "Alpha Selected": [0.05, 0.05, 0.05, 0.05],
    })

# Read uploaded file
if uploaded is not None:
    if uploaded.name.endswith(".xlsx"):
        df_raw = pd.read_excel(uploaded)  # first sheet by default
    else:
        df_raw = pd.read_csv(uploaded)
else:
    st.info("No file uploaded — using sample pre‑aggregated data.")
    df_raw = sample_aggregated()

st.subheader("Data preview")
st.dataframe(df_raw, use_container_width=True)


# ------------------------ CLEAR & CLEANUP HELPERS ------------------------
def clear_uploaded_data():
    """Clear any uploaded file and in-memory artifacts, then rerun."""
    for key in ["uploaded_file", "df_raw", "img_png_bytes"]:
        st.session_state.pop(key, None)
    # Reset the uploader widget by changing its key
    st.session_state["upload_key"] = st.session_state.get("upload_key", 0) + 1
    st.success("Uploaded data cleared from the session.")
    st.rerun()

def auto_delete_after_export(clicked: bool):
    """If a download button was clicked, clear uploaded data immediately."""
    if clicked:
        clear_uploaded_data()

# Place a manual clear button near the top (visible in both modes)
st.button("🧹 Clear uploaded data", on_click=clear_uploaded_data)


# ------------------------ Utilities ------------------------
def coerce_datetime(series: pd.Series) -> pd.Series:
    return pd.to_datetime(series, errors="coerce")

def aggregate_from_ranges(
    df_ts: pd.DataFrame,
    date_col: str,
    inf_col: str,
    bd_col: str,
    groups_spec: List[Dict[str, Tuple]]
) -> pd.DataFrame:
    """
    Aggregate raw time‑series rows into N groups based on date ranges.
    groups_spec: list of dicts like {"label": str, "start": date, "end": date, "enabled": bool}
    Returns DataFrame with columns: Group, Group Infections, Group BedDays
    """
    df = df_ts.copy()
    df[date_col] = coerce_datetime(df[date_col])
    df = df.dropna(subset=[date_col])

    out = []
    for g in groups_spec:
        if not g.get("enabled", True):
            continue
        start = pd.to_datetime(g["start"])
        end = pd.to_datetime(g["end"])
        mask = (df[date_col] >= start) & (df[date_col] <= end)
        sub = df.loc[mask]
        infections = float(sub[inf_col].sum()) if inf_col in sub.columns else np.nan
        beddays = float(sub[bd_col].sum()) if bd_col in sub.columns else np.nan
        out.append({"Group": g["label"], "Group Infections": infections, "Group BedDays": beddays})

    return pd.DataFrame(out)

def ranges_overlap(a_start: date, a_end: date, b_start: date, b_end: date) -> bool:
    return (a_start <= b_end) and (b_start <= a_end)

# parsing helpers for text inputs
def parse_date_text(text_value: str) -> date | None:
    """Parse text like 'YYYY-MM-DD' into date; returns None if invalid."""
    if not text_value:
        return None
    try:
        dt = pd.to_datetime(text_value, errors="raise")
        return dt.date()
    except Exception:
        return None

def group_range_inputs(
    group_label: str,
    default_start: date,
    default_end: date,
    key_prefix: str,
    enabled: bool,
    use_picker: bool,
    min_dt: date,
    max_dt: date,
    disabled: bool = False
) -> Tuple[str, date | None, date | None, bool]:
    """
    Render inputs for one group's label and date range.
    Returns (label, start_date, end_date, enabled_flag)
    """
    # Name field
    name = st.text_input(f"{group_label} name", value=group_label, disabled=disabled, key=f"{key_prefix}_name")
    # Enable/disable if optional
    enabled_flag = enabled
    if group_label not in ["Group 1", "Group 2"]:
        enabled_flag = st.checkbox(f"Enable {name if name else group_label}", value=enabled, key=f"{key_prefix}_enable")

    # Date inputs
    start_date, end_date = None, None
    if use_picker:
        # Date picker with a range (tuple)
        dr = st.date_input(
            f"Date range for {name if name else group_label}",
            value=(default_start, default_end),
            min_value=min_dt,
            max_value=max_dt,
            key=f"{key_prefix}_picker",
            disabled=not enabled_flag or disabled
        )
        if isinstance(dr, tuple) and len(dr) == 2:
            start_date, end_date = dr[0], dr[1]
    else:
        # Text fields with parsing/validation
        colA, colB = st.columns(2)
        with colA:
            start_text = st.text_input(
                f"Start (YYYY-MM-DD) — {name if name else group_label}",
                value=str(default_start),
                key=f"{key_prefix}_start",
                disabled=not enabled_flag or disabled,
                help="Type a date as YYYY-MM-DD (e.g., 2025-01-01)"
            )
        with colB:
            end_text = st.text_input(
                f"End (YYYY-MM-DD) — {name if name else group_label}",
                value=str(default_end),
                key=f"{key_prefix}_end",
                disabled=not enabled_flag or disabled,
                help="Type a date as YYYY-MM-DD (e.g., 2025-12-31)"
            )
        start_date = parse_date_text(start_text)
        end_date = parse_date_text(end_text)

        # Validation messages
        if enabled_flag and not disabled:
            if start_date is None:
                st.warning(f"Invalid start date for **{name or group_label}**. Please use YYYY-MM-DD.")
            if end_date is None:
                st.warning(f"Invalid end date for **{name or group_label}**. Please use YYYY-MM-DD.")
            if start_date and end_date and (start_date > end_date):
                st.warning(f"Start date is after end date for **{name or group_label}**.")
            if start_date and (start_date < min_dt or start_date > max_dt):
                st.info(f"Start date for **{name or group_label}** is outside the data range ({min_dt} → {max_dt}).")
            if end_date and (end_date < min_dt or end_date > max_dt):
                st.info(f"End date for **{name or group_label}** is outside the data range ({min_dt} → {max_dt}).")

    return name or group_label, start_date, end_date, enabled_flag


# ------------------------ Core chart logic ------------------------
def compute_and_plot(
    df: pd.DataFrame,
    alpha_input: float,
    alpha_override_flag: bool,
    show_ann: bool,
    show_ci: bool
):
    # Required columns for aggregated input
    required = ["Group", "Group Infections", "Group BedDays"]
    for col in required:
        if col not in df.columns:
            fig, ax = plt.subplots(figsize=(6, 3))
            ax.axis('off')
            ax.text(0.5, 0.5, f"Missing required column: {col}", ha='center', va='center', fontsize=12)
            plt.tight_layout()
            return fig, None, None, None, None, "n/a"

    # Determine alpha (explicit override logic)
    if (not alpha_override_flag) and ("Alpha Selected" in df.columns) and df["Alpha Selected"].notna().any():
        alpha = float(df["Alpha Selected"].iloc[0])
        alpha_source = "from data (Alpha Selected)"
    else:
        alpha = float(alpha_input)
        alpha_source = "from manual slider"

    # Clean & order
    clean_df = df[["Group", "Group Infections", "Group BedDays"]].dropna().drop_duplicates()
    labels = clean_df["Group"].astype(str).tolist()
    counts = clean_df["Group Infections"].astype(float).values
    nobs   = clean_df["Group BedDays"].astype(float).values

    # Guardrails
    mask_valid = nobs > 0
    if (~mask_valid).any() or len(clean_df) < 2:
        fig, ax = plt.subplots(figsize=(6, 3))
        ax.axis('off')
        msg = "Need ≥ 2 groups and positive bedDays in each group.\n"
        bad = [f"{labels[i]}: bedDays={int(nobs[i])}" for i in range(len(nobs)) if nobs[i] <= 0]
        if bad:
            msg += "Invalid groups → " + ", ".join(bad)
        ax.text(0.5, 0.5, msg, ha='center', va='center', fontsize=12)
        plt.tight_layout()
        return fig, None, None, None, alpha, alpha_source

    # ---------- Rates & Wilson CIs (per 1,000 BDOC) ----------
    rates = (counts / nobs) * 1000.0
    lower, upper = proportion_confint(count=counts, nobs=nobs, alpha=alpha, method='wilson')
    lower_1000, upper_1000 = lower * 1000.0, upper * 1000.0
    yerr = np.abs(np.vstack([rates - lower_1000, upper_1000 - rates]))

    # ---------- Pairwise p-value matrix ----------
    k = len(labels)
    P = np.full((k, k), np.nan)
    for i, j in itertools.combinations(range(k), 2):
        stat, pval = proportions_ztest([counts[i], counts[j]], [nobs[i], nobs[j]])
        P[i, j] = P[j, i] = pval
    np.fill_diagonal(P, 1.0)

    # ---------- Figure layout ----------
    fig = plt.figure(figsize=(14, 5.5))
    gs = fig.add_gridspec(2, 2, height_ratios=[1, 0.15], width_ratios=[3.5, 2])

    # Bars
    ax_b = fig.add_subplot(gs[0, 0])
    x = np.arange(k)
    ax_b.bar(x, rates, yerr=yerr, capsize=6, color='#5f9ea0')  # cadetblue
    ax_b.set_xticks(x)
    ax_b.set_xticklabels(labels, rotation=0, ha='center', fontsize=10)
    ax_b.set_ylabel('Rate per 1,000 BDOC')
    ax_b.set_title(f"MRSA Infection Rate (per 1000 BDOC) Time-Range Groups (α = {alpha:.3f})")

    # Per-bar labels
    for i, r in enumerate(rates):
        y_mid = r / 2.0
        if r <= 0:
            y_mid = max(upper_1000[i] * 0.05, 0.02)
        ax_b.text(
            i, y_mid, f"{r:.2f}",
            ha='center', va='center',
            fontsize=10, fontweight='bold', color='black',
            bbox=dict(facecolor='white', edgecolor='none', boxstyle='round,pad=0.25')
        )
        if show_ci:
            y_ci = upper_1000[i] * 1.01
            ci_label = f"{int((1 - alpha) * 100)}% CI [{lower_1000[i]:.2f}, {upper_1000[i]:.2f}]"
            ax_b.text(i, y_ci, ci_label, ha='center', va='bottom', fontsize=8, color='black')

    # Matrix
    ax_m = fig.add_subplot(gs[0, 1])
    base_cmap = mcolors.LinearSegmentedColormap.from_list(
        'pval_red_white_blue',
        ['#ff0000', '#ffff00', '#ffffff', '#D6EAF8', '#4682b4'],
        N=256
    )
    norm = TwoSlopeNorm(vmin=0.0, vcenter=alpha, vmax=1.0)
    im = ax_m.imshow(P, cmap=base_cmap, norm=norm)
    ax_m.set_title("Pairwise p-values (Selected Ranges)")
    ax_m.set_xlabel("Group Index")
    ax_m.set_ylabel("Group Index")
    ax_m.set_xticks(np.arange(k))
    ax_m.set_yticks(np.arange(k))
    ax_m.set_xticklabels(np.arange(1, k + 1))
    ax_m.set_yticklabels(np.arange(1, k + 1))
    cbar = fig.colorbar(im, ax=ax_m, fraction=0.04, pad=0.06)
    cbar.ax.set_ylabel('p-value', rotation=270, labelpad=12)

    # Stars and numbers when p < alpha
    if show_ann:
        for i in range(k):
            for j in range(k):
                if i != j and not np.isnan(P[i, j]):
                    star = "*" if P[i, j] < alpha else ""
                    ax_m.text(j, i, f"{P[i, j]:.3f}{star}", ha='center', va='center', fontsize=8)

    # Bottom-right key
    ax_key = fig.add_subplot(gs[1, 1]); ax_key.axis('off')
    key_items = [f"{idx + 1}: {lbl}" for idx, lbl in enumerate(labels)]
    wrapped = textwrap.fill(" • ".join(key_items), width=60)
    ax_key.text(
        0.5, 0.5, wrapped,
        ha='center', va='center', fontsize=10,
        bbox=dict(facecolor='white', alpha=0.85, boxstyle='round,pad=0.4')
    )

    plt.tight_layout()
    return fig, rates, P, clean_df, alpha, alpha_source


# ------------------------ Export helpers ------------------------
def build_rates_df(clean_df: pd.DataFrame, rates: np.ndarray) -> pd.DataFrame:
    return pd.DataFrame({
        "Group": clean_df["Group"].values,
        "Rate per 1,000 BDOC": rates
    })

def build_p_matrix_df(P: np.ndarray) -> pd.DataFrame:
    if P is None:
        return pd.DataFrame()
    k = P.shape[0]
    idx = [f"{i+1}" for i in range(k)]
    return pd.DataFrame(P, index=idx, columns=idx)


# ---------- IMPROVED PowerPoint export (with footer) ----------
def export_pptx(img_png_bytes: bytes, rates_df: pd.DataFrame, p_df: pd.DataFrame, alpha: float) -> io.BytesIO:
    """
    PPTX with:
      - Slide 1: Title + Chart (scaled)
      - Slide 2: Rates table + p-value matrix (smaller fonts)
      - Footer on each slide: 'Generated on MM-DD-YYYY', right-aligned
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()

    # Uncomment to use widescreen slides:
    # prs.slide_width = Inches(13.333)
    # prs.slide_height = Inches(7.5)

    date_str = datetime.now().strftime("%m-%d-%Y")
    footer_text = f"Generated on {date_str}"

    left_margin = Inches(0.5)
    top_margin = Inches(0.6)

    # --- Slide 1: Title + Chart ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide1.shapes.title
    title.text = f"MRSA Infection Rate (per 1,000 BDOC) — α = {alpha:.3f}"
    title.text_frame.paragraphs[0].font.size = Pt(20)

    img_stream = io.BytesIO(img_png_bytes)
    slide1.shapes.add_picture(img_stream, left_margin, top_margin + Inches(0.6), width=Inches(9.0))

    # Footer (right bottom)
    tx_width = Inches(4.0)
    tx_height = Inches(0.3)
    tx_left = prs.slide_width - (left_margin + tx_width)
    tx_top = prs.slide_height - Inches(0.4)
    footer_box = slide1.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
    fp = footer_box.text_frame.paragraphs[0]
    fp.text = footer_text
    fp.font.size = Pt(9)
    fp.alignment = PP_ALIGN.RIGHT  # ensure right-edge alignment

    # --- Slide 2: Tables ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    t2 = slide2.shapes.title
    t2.text = "Summary Tables"
    t2.text_frame.paragraphs[0].font.size = Pt(18)

    # Rates table
    rates_rows = len(rates_df) + 1
    rates_cols = 2
    rates_left = left_margin
    rates_top = top_margin + Inches(0.6)
    rates_width = Inches(4.5)
    rates_height = Inches(4.2)
    tbl_rates = slide2.shapes.add_table(rates_rows, rates_cols, rates_left, rates_top, rates_width, rates_height).table
    tbl_rates.cell(0, 0).text = "Group"
    tbl_rates.cell(0, 1).text = "Rate per 1,000 BDOC"
    for j in range(rates_cols):
        p = tbl_rates.cell(0, j).text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
    for i, row in rates_df.reset_index(drop=True).iterrows():
        tbl_rates.cell(i+1, 0).text = str(row["Group"])
        tbl_rates.cell(i+1, 1).text = f"{row['Rate per 1,000 BDOC']:.2f}"
        for j in range(rates_cols):
            tbl_rates.cell(i+1, j).text_frame.paragraphs[0].font.size = Pt(10)
    tbl_rates.columns[0].width = Inches(2.8)
    tbl_rates.columns[1].width = Inches(1.7)

    # p-value matrix table
    k = len(p_df.index)
    p_rows = k + 1
    p_cols = k + 1
    p_left = rates_left + rates_width + Inches(0.2)
    p_top = rates_top
    p_width = Inches(4.5)
    p_height = rates_height
    tbl_p = slide2.shapes.add_table(p_rows, p_cols, p_left, p_top, p_width, p_height).table
    tbl_p.cell(0, 0).text = "Row/Col"
    hdr = tbl_p.cell(0, 0).text_frame.paragraphs[0]
    hdr.font.bold = True
    hdr.font.size = Pt(11)
    for j, col_name in enumerate(p_df.columns, start=1):
        tbl_p.cell(0, j).text = str(col_name)
        h = tbl_p.cell(0, j).text_frame.paragraphs[0]
        h.font.bold = True
        h.font.size = Pt(11)
    for i, idx_label in enumerate(p_df.index, start=1):
        tbl_p.cell(i, 0).text = str(idx_label)
        rpar = tbl_p.cell(i, 0).text_frame.paragraphs[0]
        rpar.font.bold = True
        rpar.font.size = Pt(10)
        row_vals = p_df.loc[idx_label].values
        for j, val in enumerate(row_vals, start=1):
            tbl_p.cell(i, j).text = f"{val:.3f}"
            tbl_p.cell(i, j).text_frame.paragraphs[0].font.size = Pt(10)
    tbl_p.columns[0].width = Inches(1.0)
    if p_cols > 1:
        even_w = (p_width - tbl_p.columns[0].width) / (p_cols - 1)
        for c in range(1, p_cols):
            tbl_p.columns[c].width = int(even_w)

    # Footer on slide 2
    footer_box2 = slide2.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
    fp2 = footer_box2.text_frame.paragraphs[0]
    fp2.text = footer_text
    fp2.font.size = Pt(9)
    fp2.alignment = PP_ALIGN.RIGHT

    # Output buffer
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ---------- IMPROVED PDF export (landscape + footer) ----------
def export_pdf(img_png_bytes: bytes, rates_df: pd.DataFrame, p_df: pd.DataFrame, alpha: float) -> io.BytesIO:
    """
    Two-page PDF (landscape letter) via ReportLab Platypus:
      - Page 1: Title + scaled chart image
      - Page 2: Rates table + p-value matrix
      - Footer on each page: 'Generated on MM-DD-YYYY' bottom-right
    """
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors

    # Footer string
    date_str = datetime.now().strftime("%m-%d-%Y")
    footer_text = f"Generated on {date_str}"

    # PDF buffer
    out = io.BytesIO()

    # Landscape document with margins
    doc = SimpleDocTemplate(
        out,
        pagesize=landscape(letter),
        leftMargin=0.7 * inch,
        rightMargin=0.7 * inch,
        topMargin=0.6 * inch,
        bottomMargin=0.6 * inch,
    )

    # Styles
    styles = getSampleStyleSheet()
    title_style = styles["Heading2"]
    title_style.fontSize = 16
    title_style.leading = 18

    header_style = styles["Heading4"]
    header_style.fontSize = 12
    header_style.leading = 14

    # ----- Footer callback (bottom-right) -----
    def add_footer(canvas_obj, doc_obj):
        canvas_obj.setFont("Helvetica", 9)
        page_w, page_h = landscape(letter)
        # Draw so the RIGHT edge lands at (page_w - rightMargin)
        canvas_obj.drawRightString(
            page_w - doc_obj.rightMargin,
            doc_obj.bottomMargin / 2,       # vertical position above bottom edge
            footer_text
        )

    # ----- Build content -----
    story = []

    # Page 1: Title + chart
    story.append(Paragraph(f"MRSA Infection Rate (per 1,000 BDOC) — α = {alpha:.3f}", title_style))
    story.append(Spacer(1, 0.25 * inch))

    img = Image(io.BytesIO(img_png_bytes))
    # Max size for landscape letter content area
    img._restrictSize(9.5 * inch, 5.2 * inch)
    story.append(img)
    story.append(PageBreak())

    # Page 2: Tables
    story.append(Paragraph("Summary Tables", header_style))
    story.append(Spacer(1, 0.15 * inch))

    # Rates table
    rates_data = [["Group", "Rate per 1,000 BDOC"]] + [
        [str(g), f"{r:.2f}"] for g, r in zip(rates_df["Group"], rates_df["Rate per 1,000 BDOC"])
    ]
    tbl_rates = Table(rates_data, colWidths=[5.0 * inch, 2.2 * inch])
    tbl_rates.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#f0f0f0")),
        ('TEXTCOLOR', (0,0), (-1,0), colors.black),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,0), 11),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('FONTSIZE', (0,1), (-1,-1), 10),
        ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor("#fbfbfb")]),
    ]))
    story.append(tbl_rates)
    story.append(Spacer(1, 0.25 * inch))

    # p-matrix table with headers
    p_header = ["Row/Col"] + list(p_df.columns)
    p_data = [p_header] + [
        [idx] + [f"{v:.3f}" for v in p_df.loc[idx].values]
        for idx in p_df.index
    ]
    # Auto col widths: first col narrow, remaining distributed
    total_width = 9.5 * inch
    first_col_w = 1.4 * inch
    remain = total_width - first_col_w
    k = len(p_df.columns)
    col_widths = [first_col_w] + ([remain / k] * k)

    tbl_p = Table(p_data, colWidths=col_widths)
    tbl_p.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#f0f0f0")),
        ('TEXTCOLOR', (0,0), (-1,0), colors.black),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,0), 11),
        ('ALIGN', (1,1), (-1,-1), 'CENTER'),
        ('FONTSIZE', (0,1), (-1,-1), 10),
        ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor("#fbfbfb")]),
    ]))
    story.append(tbl_p)

    # Build with footer on all pages
    doc.build(story, onFirstPage=add_footer, onLaterPages=add_footer)
    out.seek(0)
    return out


# ------------------------ Modes ------------------------
alpha_effective = alpha_manual  # pass slider value; compute_and_plot decides final alpha using override flag

# --- Mode 1: Pre‑aggregated groups ---
if mode == "Pre‑aggregated groups":
    fig, rates, P, clean_df, alpha_used, alpha_source = compute_and_plot(
        df_raw.copy(), alpha_effective, alpha_override, show_annotations, show_ci_labels
    )

    # Save PNG before rendering (used for exports)
    buf_png = io.BytesIO()
    fig.savefig(buf_png, format="png", dpi=160, bbox_inches="tight")
    img_png_bytes = buf_png.getvalue()
    st.session_state["img_png_bytes"] = img_png_bytes  # store to clear later

    st.pyplot(fig, clear_figure=False)
    st.caption(f"α used: {alpha_used:.3f} — {alpha_source}")

    rates_df = None
    p_df = None
    if rates is not None and clean_df is not None:
        st.subheader("Calculated rates (per 1,000 BDOC)")
        rates_df = build_rates_df(clean_df, rates)
        st.dataframe(rates_df, use_container_width=True)

    if P is not None:
        st.subheader("Pairwise p-value matrix")
        p_df = build_p_matrix_df(P)
        st.dataframe(p_df.style.format("{:.3f}"), use_container_width=True)

    # Downloads (auto-delete after export)
    png_clicked = st.download_button(
        "Download figure as PNG",
        data=img_png_bytes,
        file_name="mrsa_rates_pvalues.png",
        mime="image/png"
    )
    auto_delete_after_export(png_clicked)

    col1, col2 = st.columns(2)
    with col1:
        if pptx_available and (rates_df is not None) and (p_df is not None):
            pptx_buf = export_pptx(img_png_bytes, rates_df, p_df, alpha_used if alpha_used is not None else alpha_effective)
            pptx_clicked = st.download_button(
                "Download PowerPoint (.pptx)",
                data=pptx_buf.getvalue(),
                file_name="mrsa_rates_pvalues.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            auto_delete_after_export(pptx_clicked)
        elif not pptx_available:
            st.info("PowerPoint export unavailable. Install: `pip install python-pptx`")
    with col2:
        if pdf_available and (rates_df is not None) and (p_df is not None):
            pdf_buf = export_pdf(img_png_bytes, rates_df, p_df, alpha_used if alpha_used is not None else alpha_effective)
            pdf_clicked = st.download_button(
                "Download PDF (.pdf)",
                data=pdf_buf.getvalue(),
                file_name="mrsa_rates_pvalues.pdf",
                mime="application/pdf"
            )
            auto_delete_after_export(pdf_clicked)
        elif not pdf_available:
            st.info("PDF export unavailable. Install: `pip install reportlab`")

# --- Mode 2: Custom date ranges (up to 4 groups) ---
else:
    st.markdown("**Map your columns** (for raw time‑series data)")
    # Auto-suggestions
    date_candidates = [c for c in df_raw.columns if "date" in c.lower()]
    inf_candidates = [c for c in df_raw.columns if "infect" in c.lower()] + [c for c in df_raw.columns if "case" in c.lower()]
    bd_candidates  = [c for c in df_raw.columns if "bed" in c.lower()] + [c for c in df_raw.columns if "bdoc" in c.lower()]

    date_col = st.selectbox("Date column", options=date_candidates if date_candidates else list(df_raw.columns))
    inf_col  = st.selectbox("Infections column", options=inf_candidates if inf_candidates else list(df_raw.columns))
    bd_col   = st.selectbox("BedDays column", options=bd_candidates if bd_candidates else list(df_raw.columns))

    # Determine overall min/max dates to seed the inputs
    df_dates = df_raw.copy()
    df_dates[date_col] = coerce_datetime(df_dates[date_col])
    df_dates = df_dates.dropna(subset=[date_col])
    if len(df_dates) == 0:
        st.warning("No valid dates found in the selected Date column.")
        st.stop()

    min_dt = pd.to_datetime(df_dates[date_col].min()).date()
    max_dt = pd.to_datetime(df_dates[date_col].max()).date()
    st.write(f"Data date range detected: **{min_dt} → {max_dt}**")

    # Up to 4 groups with custom names and toggle to delete/disable
    groups_spec = []

    st.subheader("Groups")

    # Group 1 (always enabled)
    g1_label, g1_start, g1_end, _g1_enabled = group_range_inputs(
        "Group 1", min_dt, max_dt, "g1", True, use_date_picker, min_dt, max_dt, disabled=False
    )
    groups_spec.append({"label": g1_label, "start": g1_start, "end": g1_end, "enabled": True})

    # Group 2 (always enabled)
    g2_label, g2_start, g2_end, _g2_enabled = group_range_inputs(
        "Group 2", min_dt, max_dt, "g2", True, use_date_picker, min_dt, max_dt, disabled=False
    )
    groups_spec.append({"label": g2_label, "start": g2_start, "end": g2_end, "enabled": True})

    # Group 3 (optional)
    g3_label, g3_start, g3_end, g3_enabled = group_range_inputs(
        "Group 3", min_dt, max_dt, "g3", False, use_date_picker, min_dt, max_dt, disabled=False
    )
    groups_spec.append({"label": g3_label, "start": g3_start, "end": g3_end, "enabled": bool(g3_enabled)})

    # Group 4 (optional)
    g4_label, g4_start, g4_end, g4_enabled = group_range_inputs(
        "Group 4", min_dt, max_dt, "g4", False, use_date_picker, min_dt, max_dt, disabled=False
    )
    groups_spec.append({"label": g4_label, "start": g4_start, "end": g4_end, "enabled": bool(g4_enabled)})

    # Validate enabled groups have valid dates
    enabled_groups = [g for g in groups_spec if g["enabled"]]
    invalid_groups = []
    for g in enabled_groups:
        if (g["start"] is None) or (g["end"] is None) or (g["start"] > g["end"]):
            invalid_groups.append(g["label"])
    if invalid_groups:
        st.error("Invalid date ranges for: " + ", ".join(invalid_groups))
        st.stop()

    # Overlap warnings (pairwise across enabled groups)
    for i in range(len(enabled_groups)):
        for j in range(i + 1, len(enabled_groups)):
            gi, gj = enabled_groups[i], enabled_groups[j]
            if ranges_overlap(gi["start"], gi["end"], gj["start"], gj["end"]):
                st.warning(f"Date ranges overlap: **{gi['label']}** and **{gj['label']}**. Overlapping periods will count in both groups.")

    # Require at least 2 enabled groups
    if len(enabled_groups) < 2:
        st.info("Please enable at least **two** groups.")
        st.stop()

    # Aggregate into selected groups
    agg_df = aggregate_from_ranges(df_raw, date_col, inf_col, bd_col, groups_spec)

    # Compute and plot
    fig, rates, P, clean_df, alpha_used, alpha_source = compute_and_plot(
        agg_df.copy(), alpha_effective, alpha_override, show_annotations, show_ci_labels
    )

    # Save PNG before rendering (used for exports)
    buf_png = io.BytesIO()
    fig.savefig(buf_png, format="png", dpi=160, bbox_inches="tight")
    img_png_bytes = buf_png.getvalue()
    st.session_state["img_png_bytes"] = img_png_bytes

    st.pyplot(fig, clear_figure=False)
    st.caption(f"α used: {alpha_used:.3f} — {alpha_source}")

    # Tables
    rates_df = None
    p_df = None
    if rates is not None and clean_df is not None:
        st.subheader("Calculated rates (per 1,000 BDOC)")
        rates_df = build_rates_df(clean_df, rates)
        st.dataframe(rates_df, use_container_width=True)

    if P is not None:
        st.subheader("Pairwise p-value matrix")
        p_df = build_p_matrix_df(P)
        st.dataframe(p_df.style.format("{:.3f}"), use_container_width=True)

    # Downloads (auto-delete after export)
    png_clicked = st.download_button(
        "Download figure as PNG",
        data=img_png_bytes,
        file_name="mrsa_rates_pvalues.png",
        mime="image/png"
    )
    auto_delete_after_export(png_clicked)

    col1, col2 = st.columns(2)
    with col1:
        if pptx_available and (rates_df is not None) and (p_df is not None):
            pptx_buf = export_pptx(img_png_bytes, rates_df, p_df, alpha_used if alpha_used is not None else alpha_effective)
            pptx_clicked = st.download_button(
                "Download PowerPoint (.pptx)",
                data=pptx_buf.getvalue(),
                file_name="mrsa_rates_pvalues.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            auto_delete_after_export(pptx_clicked)
        elif not pptx_available:
            st.info("PowerPoint export unavailable. Install: `pip install python-pptx`")
    with col2:
        if pdf_available and (rates_df is not None) and (p_df is not None):
            pdf_buf = export_pdf(img_png_bytes, rates_df, p_df, alpha_used if alpha_used is not None else alpha_effective)
            pdf_clicked = st.download_button(
                "Download PDF (.pdf)",
                data=pdf_buf.getvalue(),
                file_name="mrsa_rates_pvalues.pdf",
                mime="application/pdf"
            )
            auto_delete_after_export(pdf_clicked)
        elif not pdf_available:
            st.info("PDF export unavailable. Install: `pip install reportlab`")
